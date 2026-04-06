"""Media-Pipeline: Verarbeitung von Bildern, Audio und Dokumenten.

MCP-Tools fuer multimodale Medienverarbeitung -- vollstaendig lokal.

Tools:
  - media_transcribe_audio: Audio → Text (Whisper)
  - media_analyze_image: Bild → Beschreibung (multimodales LLM via Ollama)
  - media_extract_text: PDF/DOCX/TXT → Text
  - media_convert_audio: Audio-Formatkonvertierung (ffmpeg)
  - media_image_resize: Bildgroesse aendern (Pillow)
  - media_tts: Text → Sprache (Piper/eSpeak)

Alle Tools arbeiten mit lokalen Dateipfaden -- keine Cloud-Uploads.
"""

from __future__ import annotations

import asyncio
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from jarvis.i18n import t
from jarvis.utils.logging import get_logger

log = get_logger(__name__)

# Maximale Textlaenge fuer LLM-Kontext
_DEFAULT_MAX_EXTRACT_LENGTH = 15_000

# Maximale Bilddateigroesse fuer Base64-Encoding (10 MB)
_DEFAULT_MAX_IMAGE_FILE_SIZE = 10_485_760

# Maximale Dateigroessen fuer Security-Limits
_DEFAULT_MAX_EXTRACT_FILE_SIZE = 52_428_800  # 50 MB für Dokument-Extraktion
_DEFAULT_MAX_AUDIO_FILE_SIZE = 104_857_600  # 100 MB für Audio-Transkription

# Maximale Bilddimensionen und Standard-Resize-Werte
_DEFAULT_MAX_IMAGE_DIMENSION = 8192
_DEFAULT_MAX_WIDTH = 1024
_DEFAULT_MAX_HEIGHT = 1024

# Backward compatibility aliases
MAX_EXTRACT_LENGTH = _DEFAULT_MAX_EXTRACT_LENGTH
MAX_IMAGE_FILE_SIZE = _DEFAULT_MAX_IMAGE_FILE_SIZE
MAX_EXTRACT_FILE_SIZE = _DEFAULT_MAX_EXTRACT_FILE_SIZE
MAX_AUDIO_FILE_SIZE = _DEFAULT_MAX_AUDIO_FILE_SIZE

# Standard-Modelle und -Stimmen (Fallbacks, wenn kein Config verfuegbar)
_DEFAULT_VISION_MODEL = "openbmb/minicpm-v4.5"
_DEFAULT_VISION_MODEL_DETAIL = "qwen3-vl:32b"
DEFAULT_IMAGE_PROMPT = "Beschreibe dieses Bild detailliert auf Deutsch."
DEFAULT_PIPER_VOICE = "de_DE-thorsten-high"

__all__ = [
    "MEDIA_TOOL_SCHEMAS",
    "MediaPipeline",
    "MediaResult",
    "register_media_tools",
]


@dataclass
class MediaResult:
    """Einheitliches Ergebnis aller Media-Operationen."""

    success: bool = True
    text: str = ""
    output_path: str | None = None
    metadata: dict[str, Any] | None = None
    error: str | None = None


class MediaPipeline:
    """Zentrale Klasse fuer Medienverarbeitung.

    Alle Methoden sind async und nutzen run_in_executor
    fuer CPU-intensive Operationen (Whisper, Pillow, etc.).
    """

    def __init__(self, workspace_dir: Path | None = None, config: Any = None) -> None:
        self._workspace = workspace_dir or Path.home() / ".jarvis" / "workspace" / "media"
        self._workspace.mkdir(parents=True, exist_ok=True)
        # LLM + Vault injection (gesetzt via _set_llm_fn / _set_vault)
        self._llm_fn: Any = None
        self._llm_model: str = ""
        self._vault: Any = None
        self._config = config
        # Lazy-initialised template manager
        self._template_manager: Any = None

        # Konfigurierbare Limits (aus config.media.* mit Fallback auf Defaults)
        _media = getattr(config, "media", None)
        self._max_extract_length: int = getattr(
            _media, "max_extract_length", _DEFAULT_MAX_EXTRACT_LENGTH
        )
        self._max_image_file_size: int = getattr(
            _media, "max_image_file_size", _DEFAULT_MAX_IMAGE_FILE_SIZE
        )
        self._max_extract_file_size: int = getattr(
            _media, "max_extract_file_size", _DEFAULT_MAX_EXTRACT_FILE_SIZE
        )
        self._max_audio_file_size: int = getattr(
            _media, "max_audio_file_size", _DEFAULT_MAX_AUDIO_FILE_SIZE
        )
        self._max_image_dimension: int = getattr(
            _media, "max_image_dimension", _DEFAULT_MAX_IMAGE_DIMENSION
        )
        self._default_max_width: int = getattr(_media, "default_max_width", _DEFAULT_MAX_WIDTH)
        self._default_max_height: int = getattr(_media, "default_max_height", _DEFAULT_MAX_HEIGHT)

    def _set_llm_fn(self, llm_fn: Any, model_name: str = "") -> None:
        """Injiziert eine LLM-Funktion fuer Dokument-Analyse.

        Args:
            llm_fn: Async-Funktion mit Signatur (prompt: str, model: str) -> str
            model_name: Name des zu verwendenden Modells.
        """
        self._llm_fn = llm_fn
        self._llm_model = model_name

    def _set_vault(self, vault: Any) -> None:
        """Injiziert VaultTools-Referenz fuer optionales Speichern.

        Args:
            vault: VaultTools-Instanz mit vault_save-Methode.
        """
        self._vault = vault

    def _validate_input_path(self, file_path: str) -> Path | None:
        """Validates input file path against path traversal and workspace confinement.

        Returns resolved Path or None if invalid.
        """
        try:
            path = Path(file_path).expanduser().resolve()
        except (ValueError, OSError):
            return None
        if not path.exists():
            return None
        # Workspace confinement: path must be inside workspace or home/.jarvis
        jarvis_home = Path.home() / ".jarvis"
        try:
            path.relative_to(self._workspace)
        except ValueError:
            try:
                path.relative_to(jarvis_home)
            except ValueError:
                log.warning(
                    "media_path_outside_workspace", path=str(path), workspace=str(self._workspace)
                )
                return None
        return path

    # ========================================================================
    # Audio → Text (Whisper STT)
    # ========================================================================

    async def transcribe_audio(
        self,
        audio_path: str,
        *,
        language: str = "de",
        model: str = "base",
    ) -> MediaResult:
        """Transkribiert eine Audiodatei zu Text.

        Unterstuetzt: WAV, MP3, OGG, FLAC, M4A, WEBM
        Backend: faster-whisper (lokal, GPU-beschleunigt)

        Args:
            audio_path: Pfad zur Audiodatei.
            language: Sprache (ISO-Code, z.B. 'de', 'en').
            model: Whisper-Modell ('tiny', 'base', 'small', 'medium', 'large-v3').
        """
        path = self._validate_input_path(audio_path)
        if path is None:
            return MediaResult(
                success=False, error=t("media.file_not_found", path=audio_path)
            )

        file_size = path.stat().st_size
        if file_size > self._max_audio_file_size:
            return MediaResult(
                success=False,
                error=t(
                    "media.audio_too_large",
                    size_mb=f"{file_size / 1_048_576:.1f}",
                    max_mb=self._max_audio_file_size // 1_048_576,
                ),
            )

        try:
            from faster_whisper import WhisperModel

            loop = asyncio.get_running_loop()

            def _transcribe() -> str:
                m = WhisperModel(model, device="auto", compute_type="int8")
                segments, info = m.transcribe(str(path), language=language, vad_filter=True)
                text = " ".join(seg.text.strip() for seg in segments)
                return text

            text = await loop.run_in_executor(None, _transcribe)

            if not text.strip():
                return MediaResult(success=True, text=t("media.no_speech_detected"))

            log.info("audio_transcribed", path=audio_path, length=len(text))
            return MediaResult(
                success=True,
                text=text,
                metadata={"language": language, "model": model, "source": audio_path},
            )

        except ImportError:
            return MediaResult(
                success=False,
                error=t("media.whisper_not_installed"),
            )
        except Exception as exc:
            log.error("transcribe_failed", path=audio_path, error=str(exc))
            return MediaResult(success=False, error=t("media.transcription_failed", error=exc))

    # ========================================================================
    # Bild → Beschreibung (multimodales LLM)
    # ========================================================================

    async def analyze_image(
        self,
        image_path: str,
        *,
        prompt: str = DEFAULT_IMAGE_PROMPT,
        model: str = _DEFAULT_VISION_MODEL,
        ollama_url: str = "http://localhost:11434",
        openai_api_key: str = "",
        openai_base_url: str = "https://api.openai.com/v1",
    ) -> MediaResult:
        """Analysiert ein Bild mit einem multimodalen LLM.

        Unterstuetzt: JPG, PNG, GIF, WEBP, BMP
        Backend: Automatisch erkannt — OpenAI fuer gpt-*/o*-Modelle, sonst Ollama.

        Args:
            image_path: Pfad zum Bild.
            prompt: Analyseanweisung fuer das LLM.
            model: Vision-Modell (OpenAI oder Ollama).
            ollama_url: Ollama API-URL.
            openai_api_key: OpenAI API-Key (noetig fuer gpt-* Modelle).
            openai_base_url: OpenAI base URL.
        """
        import base64

        path = self._validate_input_path(image_path)
        if path is None:
            return MediaResult(
                success=False, error=t("media.image_not_found", path=image_path)
            )

        suffix = path.suffix.lower()
        if suffix not in (".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp"):
            return MediaResult(success=False, error=t("media.unsupported_image", suffix=suffix))

        try:
            import httpx

            # Dateigroesse pruefen
            file_size = path.stat().st_size
            if file_size > self._max_image_file_size:
                return MediaResult(
                    success=False,
                    error=t(
                        "media.image_too_large",
                        size_mb=f"{file_size / 1_048_576:.1f}",
                        max_mb=self._max_image_file_size // 1_048_576,
                    ),
                )

            # Bild als Base64 laden
            image_data = base64.b64encode(path.read_bytes()).decode("utf-8")

            # Backend-Erkennung: gpt-* und o*-Modelle → OpenAI, sonst Ollama
            _is_openai = openai_api_key and (
                model.startswith("gpt-")
                or model.startswith("o1")
                or model.startswith("o3")
                or model.startswith("o4")
            )

            async with httpx.AsyncClient(timeout=180.0, trust_env=False) as client:
                if _is_openai:
                    # OpenAI Chat Completions mit Vision
                    suffix = path.suffix.lower().lstrip(".")
                    mime = {
                        "jpg": "jpeg",
                        "jpeg": "jpeg",
                        "png": "png",
                        "gif": "gif",
                        "webp": "webp",
                        "bmp": "bmp",
                    }.get(suffix, "jpeg")
                    resp = await client.post(
                        f"{openai_base_url}/chat/completions",
                        headers={"Authorization": f"Bearer {openai_api_key}"},
                        json={
                            "model": model,
                            "messages": [
                                {
                                    "role": "user",
                                    "content": [
                                        {"type": "text", "text": prompt},
                                        {
                                            "type": "image_url",
                                            "image_url": {
                                                "url": f"data:image/{mime};base64,{image_data}"
                                            },
                                        },
                                    ],
                                }
                            ],
                            "max_completion_tokens": 2000,
                        },
                    )
                    if resp.status_code != 200:
                        return MediaResult(
                            success=False,
                            error=f"OpenAI HTTP {resp.status_code}: {resp.text[:300]}",
                        )
                    data = resp.json()
                    description = (
                        (data.get("choices") or [{}])[0].get("message", {}).get("content", "")
                    )
                else:
                    # Ollama /api/chat mit images-Array
                    resp = await client.post(
                        f"{ollama_url}/api/chat",
                        json={
                            "model": model,
                            "messages": [
                                {
                                    "role": "user",
                                    "content": prompt,
                                    "images": [image_data],
                                }
                            ],
                            "stream": False,
                        },
                    )
                    if resp.status_code != 200:
                        return MediaResult(
                            success=False,
                            error=f"Ollama HTTP {resp.status_code}: {resp.text[:300]}",
                        )
                    data = resp.json()
                    description = data.get("message", {}).get("content", "")

                log.info("image_analyzed", path=image_path, model=model)
                return MediaResult(
                    success=True,
                    text=description,
                    metadata={
                        "model": model,
                        "source": image_path,
                        "image_size": path.stat().st_size,
                    },
                )

        except ImportError:
            return MediaResult(success=False, error=t("media.httpx_not_installed"))
        except Exception as exc:
            log.error("image_analysis_failed", path=image_path, error=str(exc))
            return MediaResult(success=False, error=t("media.image_analysis_failed", error=exc))

    # ========================================================================
    # Dokument → Text
    # ========================================================================

    async def extract_text(self, file_path: str) -> MediaResult:
        """Extrahiert Text aus verschiedenen Dokumentformaten.

        Unterstuetzt: PDF, DOCX, TXT, MD, HTML, CSV, JSON, XML

        Args:
            file_path: Pfad zum Dokument.
        """
        path = self._validate_input_path(file_path)
        if path is None:
            return MediaResult(
                success=False, error=t("media.file_not_found", path=file_path)
            )

        file_size = path.stat().st_size
        if file_size > self._max_extract_file_size:
            return MediaResult(
                success=False,
                error=t(
                    "media.file_too_large_mb",
                    size_mb=f"{file_size / 1_048_576:.1f}",
                    max_mb=self._max_extract_file_size // 1_048_576,
                ),
            )

        suffix = path.suffix.lower()
        loop = asyncio.get_running_loop()

        try:
            if suffix == ".pdf":
                text = await loop.run_in_executor(None, self._extract_pdf, path)
            elif suffix == ".docx":
                text = await loop.run_in_executor(None, self._extract_docx, path)
            elif suffix in (".txt", ".md", ".csv", ".json", ".xml", ".yaml", ".yml", ".log"):
                text = path.read_text(encoding="utf-8", errors="replace")
            elif suffix in (".html", ".htm"):
                text = await loop.run_in_executor(None, self._extract_html, path)
            elif suffix == ".pptx":
                text = await loop.run_in_executor(None, self._extract_ppt_text, path)
            else:
                return MediaResult(
                    success=False,
                    error=t("media.unsupported_format_extended", suffix=suffix),
                )

            original_length = len(text)
            if len(text) > self._max_extract_length:
                text = (
                    text[: self._max_extract_length]
                    + "\n\n" + t("media.text_truncated", length=original_length)
                )

            log.info("text_extracted", path=file_path, length=len(text), format=suffix)
            return MediaResult(
                success=True,
                text=text,
                metadata={"source": file_path, "format": suffix, "original_length": original_length},
            )

        except Exception as exc:
            log.error("text_extraction_failed", path=file_path, error=str(exc))
            return MediaResult(success=False, error=t("media.text_extraction_failed", error=exc))

    def _extract_pdf(self, path: Path) -> str:
        """PDF-Textextraktion mit pymupdf oder pdfplumber."""
        # Versuch 1: PyMuPDF (schnell)
        try:
            import fitz  # pymupdf

            doc = fitz.open(str(path))
            pages = []
            for page in doc:
                pages.append(page.get_text())
            doc.close()
            return "\n\n".join(pages)
        except ImportError:
            pass

        # Versuch 2: pdfplumber
        try:
            import pdfplumber

            with pdfplumber.open(str(path)) as pdf:
                pages = [page.extract_text() or "" for page in pdf.pages]
                return "\n\n".join(pages)
        except ImportError:
            pass

        raise ImportError(t("media.pdf_reader_required"))

    def _extract_docx(self, path: Path) -> str:
        """DOCX-Textextraktion mit python-docx."""
        try:
            from docx import Document

            doc = Document(str(path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except ImportError:
            raise ImportError("python-docx nicht installiert. pip install python-docx") from None

    def _extract_html(self, path: Path) -> str:
        """HTML-Textextraktion (einfach, ohne BeautifulSoup-Pflicht)."""
        import re

        if path.stat().st_size > self._max_extract_file_size:
            raise ValueError(
                f"HTML-Datei zu gross ({path.stat().st_size / 1_048_576:.1f} MB, "
                f"max {self._max_extract_file_size // 1_048_576} MB)"
            )
        html = path.read_text(encoding="utf-8", errors="replace")
        # Script/Style-Tags entfernen
        html = re.sub(r"<(script|style)[^>]*>.*?</\1>", "", html, flags=re.DOTALL | re.IGNORECASE)
        # HTML-Tags entfernen
        text = re.sub(r"<[^>]+>", " ", html)
        # Mehrfach-Whitespace normalisieren
        text = re.sub(r"\s+", " ", text).strip()
        return text

    def _extract_ppt_text(self, path: Path) -> str:
        """PPTX-Textextraktion (einfach, nur Text)."""
        try:
            from pptx import Presentation

            prs = Presentation(str(path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        texts.append(shape.text_frame.text)
            return "\n\n".join(_txt for _txt in texts if _txt.strip())
        except ImportError:
            raise ImportError("python-pptx nicht installiert. pip install python-pptx") from None

    # ========================================================================
    # Strukturierte Dokument-Lese-Tools (PDF / PPT / DOCX)
    # ========================================================================

    async def read_pdf(
        self,
        file_path: str,
        *,
        extract_images: bool = False,
        extract_tables: bool = False,
        pages: str = "",
    ) -> MediaResult:
        """Liest ein PDF strukturiert mit Seiten, Metadaten, Bildern, Tabellen."""
        path = self._validate_input_path(file_path)
        if path is None:
            return MediaResult(
                success=False, error=t("media.file_not_found", path=file_path)
            )

        file_size = path.stat().st_size
        if file_size > self._max_extract_file_size:
            return MediaResult(
                success=False,
                error=t(
                    "media.file_too_large_mb",
                    size_mb=f"{file_size / 1_048_576:.1f}",
                    max_mb=self._max_extract_file_size // 1_048_576,
                ),
            )

        loop = asyncio.get_running_loop()
        try:
            result = await loop.run_in_executor(
                None,
                self._read_pdf_structured,
                path,
                extract_images,
                extract_tables,
                pages,
            )
            return result
        except Exception as exc:
            log.error("read_pdf_failed", path=file_path, error=str(exc))
            return MediaResult(success=False, error=t("media.pdf_read_failed", error=exc))

    def _read_pdf_structured(
        self,
        path: Path,
        extract_images: bool,
        extract_tables: bool,
        pages: str,
    ) -> MediaResult:
        """Synchrone strukturierte PDF-Extraktion."""
        try:
            import fitz  # pymupdf
        except ImportError:
            raise ImportError("pymupdf nicht installiert. pip install pymupdf") from None

        doc = fitz.open(str(path))
        total_pages = len(doc)

        # Seitenbereich parsen (1-basiert -> 0-basiert)
        if pages.strip():
            if "-" in pages:
                parts = pages.split("-", 1)
                start = max(int(parts[0]) - 1, 0)
                end = min(int(parts[1]), total_pages)
                page_indices = list(range(start, end))
            else:
                idx = int(pages) - 1
                page_indices = [idx] if 0 <= idx < total_pages else []
        else:
            page_indices = list(range(total_pages))

        text_parts: list[str] = []
        image_paths: list[str] = []
        tables_md: list[str] = []

        for i in page_indices:
            page = doc[i]
            page_text = page.get_text()
            text_parts.append(f"--- Seite {i + 1} ---\n{page_text}")

            if extract_images:
                img_dir = self._workspace / "pdf_images"
                img_dir.mkdir(parents=True, exist_ok=True)
                for img_idx, img_info in enumerate(page.get_images(full=True)):
                    xref = img_info[0]
                    try:
                        base_image = doc.extract_image(xref)
                        ext = base_image.get("ext", "png")
                        img_path = img_dir / f"page{i + 1}_img{img_idx + 1}.{ext}"
                        img_path.write_bytes(base_image["image"])
                        image_paths.append(str(img_path))
                    except Exception:
                        pass

            if extract_tables:
                try:
                    found_tables = page.find_tables()
                    for t_idx, table in enumerate(found_tables):
                        rows = table.extract()
                        if rows:
                            md_rows = [
                                "| " + " | ".join(str(c) if c else "" for c in row) + " |"
                                for row in rows
                            ]
                            if len(md_rows) > 1:
                                header_sep = "| " + " | ".join("---" for _ in rows[0]) + " |"
                                md_rows.insert(1, header_sep)
                            tables_md.append(
                                f"Tabelle (Seite {i + 1}, #{t_idx + 1}):\n" + "\n".join(md_rows)
                            )
                except Exception:
                    pass

        # Metadaten
        meta = doc.metadata or {}
        metadata = {
            "page_count": total_pages,
            "pages_extracted": len(page_indices),
            "title": meta.get("title", ""),
            "author": meta.get("author", ""),
            "subject": meta.get("subject", ""),
            "creator": meta.get("creator", ""),
            "creation_date": meta.get("creationDate", ""),
        }
        doc.close()

        full_text = "\n\n".join(text_parts)
        if tables_md:
            full_text += "\n\n" + "\n\n".join(tables_md)
        if image_paths:
            full_text += "\n\nExtrahierte Bilder:\n" + "\n".join(image_paths)
            metadata["images"] = image_paths
        if tables_md:
            metadata["tables_count"] = len(tables_md)

        log.info(
            "read_pdf_done",
            path=str(path),
            pages=len(page_indices),
            images=len(image_paths),
            tables=len(tables_md),
        )
        return MediaResult(success=True, text=full_text, metadata=metadata)

    async def read_ppt(
        self,
        file_path: str,
        *,
        extract_images: bool = False,
    ) -> MediaResult:
        """Liest eine PowerPoint-Praesentation strukturiert."""
        path = self._validate_input_path(file_path)
        if path is None:
            return MediaResult(
                success=False, error=t("media.file_not_found", path=file_path)
            )

        file_size = path.stat().st_size
        if file_size > self._max_extract_file_size:
            return MediaResult(
                success=False,
                error=t(
                    "media.file_too_large_mb",
                    size_mb=f"{file_size / 1_048_576:.1f}",
                    max_mb=self._max_extract_file_size // 1_048_576,
                ),
            )

        loop = asyncio.get_running_loop()
        try:
            result = await loop.run_in_executor(
                None, self._read_ppt_structured, path, extract_images
            )
            return result
        except Exception as exc:
            log.error("read_ppt_failed", path=file_path, error=str(exc))
            return MediaResult(success=False, error=t("media.ppt_read_failed", error=exc))

    def _read_ppt_structured(self, path: Path, extract_images: bool) -> MediaResult:
        """Synchrone strukturierte PPTX-Extraktion."""
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            raise ImportError("python-pptx nicht installiert. pip install python-pptx") from None

        prs = Presentation(str(path))
        slide_texts: list[str] = []
        image_paths: list[str] = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            title = ""
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text

            content_parts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        para_parts: list[str] = []
                        for run in paragraph.runs:
                            text = run.text
                            if run.font.bold:
                                text = f"**{text}**"
                            if run.font.italic:
                                text = f"*{text}*"
                            para_parts.append(text)
                        line = "".join(para_parts)
                        if line.strip():
                            content_parts.append(line)

                if extract_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_dir = self._workspace / "ppt_images"
                    img_dir.mkdir(parents=True, exist_ok=True)
                    try:
                        image = shape.image
                        ext = image.content_type.split("/")[-1] if image.content_type else "png"
                        img_path = img_dir / f"slide{slide_idx}_{shape.shape_id}.{ext}"
                        img_path.write_bytes(image.blob)
                        image_paths.append(str(img_path))
                    except Exception:
                        pass

            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()

            slide_text = f"--- Folie {slide_idx}: {title} ---\n"
            slide_text += "\n".join(content_parts)
            if notes:
                slide_text += f"\n\nNotizen: {notes}"
            slide_texts.append(slide_text)

        full_text = "\n\n".join(slide_texts)
        metadata: dict[str, Any] = {"slide_count": len(prs.slides)}
        if image_paths:
            full_text += "\n\nExtrahierte Bilder:\n" + "\n".join(image_paths)
            metadata["images"] = image_paths

        log.info(
            "read_ppt_done",
            path=str(path),
            slides=len(prs.slides),
            images=len(image_paths),
        )
        return MediaResult(success=True, text=full_text, metadata=metadata)

    async def read_docx(
        self,
        file_path: str,
        *,
        extract_images: bool = False,
        extract_tables: bool = True,
    ) -> MediaResult:
        """Liest ein Word-Dokument strukturiert."""
        path = self._validate_input_path(file_path)
        if path is None:
            return MediaResult(
                success=False, error=t("media.file_not_found", path=file_path)
            )

        file_size = path.stat().st_size
        if file_size > self._max_extract_file_size:
            return MediaResult(
                success=False,
                error=t(
                    "media.file_too_large_mb",
                    size_mb=f"{file_size / 1_048_576:.1f}",
                    max_mb=self._max_extract_file_size // 1_048_576,
                ),
            )

        loop = asyncio.get_running_loop()
        try:
            result = await loop.run_in_executor(
                None, self._read_docx_structured, path, extract_images, extract_tables
            )
            return result
        except Exception as exc:
            log.error("read_docx_failed", path=file_path, error=str(exc))
            return MediaResult(success=False, error=t("media.docx_read_failed", error=exc))

    def _read_docx_structured(
        self, path: Path, extract_images: bool, extract_tables: bool
    ) -> MediaResult:
        """Synchrone strukturierte DOCX-Extraktion."""
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx nicht installiert. pip install python-docx") from None

        doc = Document(str(path))
        text_parts: list[str] = []

        # Absaetze mit Heading-Erkennung und Formatierung
        for paragraph in doc.paragraphs:
            style_name = paragraph.style.name if paragraph.style else ""
            if style_name.startswith("Heading"):
                try:
                    level = int(style_name.replace("Heading", "").strip())
                except ValueError:
                    level = 1
                text_parts.append("#" * level + " " + paragraph.text)
            else:
                # Runs mit Formatierung
                para_parts: list[str] = []
                for run in paragraph.runs:
                    text = run.text
                    if run.bold:
                        text = f"**{text}**"
                    if run.italic:
                        text = f"*{text}*"
                    para_parts.append(text)
                line = "".join(para_parts)
                if line.strip():
                    text_parts.append(line)

        # Tabellen als Markdown
        tables_md: list[str] = []
        if extract_tables:
            for t_idx, table in enumerate(doc.tables):
                rows: list[list[str]] = []
                for row in table.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                if rows:
                    md_rows = [
                        "| " + " | ".join(rows[0]) + " |",
                        "| " + " | ".join("---" for _ in rows[0]) + " |",
                    ]
                    for row in rows[1:]:
                        md_rows.append("| " + " | ".join(row) + " |")
                    tables_md.append(f"Tabelle #{t_idx + 1}:\n" + "\n".join(md_rows))

        # Bilder extrahieren
        image_paths: list[str] = []
        if extract_images:
            img_dir = self._workspace / "docx_images"
            img_dir.mkdir(parents=True, exist_ok=True)
            img_idx = 0
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        img_idx += 1
                        blob = rel.target_part.blob
                        content_type = rel.target_part.content_type or "image/png"
                        ext = content_type.split("/")[-1]
                        img_path = img_dir / f"image_{img_idx}.{ext}"
                        img_path.write_bytes(blob)
                        image_paths.append(str(img_path))
                    except Exception:
                        pass

        # Metadaten
        props = doc.core_properties
        metadata: dict[str, Any] = {
            "author": props.author or "",
            "title": props.title or "",
            "created": str(props.created) if props.created else "",
            "modified": str(props.modified) if props.modified else "",
            "paragraph_count": len(doc.paragraphs),
            "table_count": len(doc.tables),
        }

        full_text = "\n\n".join(text_parts)
        if tables_md:
            full_text += "\n\n" + "\n\n".join(tables_md)
        if image_paths:
            full_text += "\n\nExtrahierte Bilder:\n" + "\n".join(image_paths)
            metadata["images"] = image_paths

        log.info(
            "read_docx_done",
            path=str(path),
            paragraphs=len(doc.paragraphs),
            tables=len(tables_md),
            images=len(image_paths),
        )
        return MediaResult(success=True, text=full_text, metadata=metadata)

    # ========================================================================
    # Excel-Lesen (.xlsx)
    # ========================================================================

    async def read_xlsx(
        self,
        file_path: str,
        *,
        sheet_name: str = "",
        max_rows: int = 500,
    ) -> MediaResult:
        """Reads an Excel file (.xlsx) and returns data as Markdown tables.

        Args:
            file_path: Path to the .xlsx file.
            sheet_name: Specific sheet to read (empty = all sheets).
            max_rows: Maximum rows per sheet (default 500).
        """
        path = Path(file_path).expanduser().resolve()
        if not path.exists():
            return MediaResult(success=False, error=f"File not found: {path}")
        if path.suffix.lower() not in (".xlsx", ".xls"):
            return MediaResult(success=False, error=f"Not an Excel file: {path.suffix}")

        loop = asyncio.get_running_loop()
        try:
            return await loop.run_in_executor(
                None, self._read_xlsx_structured, path, sheet_name, max_rows
            )
        except Exception as exc:
            log.error("read_xlsx_failed", path=file_path, error=str(exc))
            return MediaResult(success=False, error=str(exc))

    def _read_xlsx_structured(self, path: Path, sheet_name: str, max_rows: int) -> MediaResult:
        """Sync worker for read_xlsx."""
        try:
            from openpyxl import load_workbook
        except ImportError:
            return MediaResult(
                success=False,
                error="openpyxl not installed. Run: pip install openpyxl",
            )

        wb = load_workbook(path, read_only=True, data_only=True)
        sheets = [sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.sheetnames

        parts = []
        total_rows = 0
        for sname in sheets:
            ws = wb[sname]
            parts.append(f"## Sheet: {sname}")

            rows = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= max_rows:
                    parts.append(f"... ({ws.max_row - max_rows} more rows truncated)")
                    break
                cells = [str(c) if c is not None else "" for c in row]
                rows.append(cells)
                total_rows += 1

            if rows:
                # Markdown table
                header = rows[0]
                parts.append("| " + " | ".join(header) + " |")
                parts.append("| " + " | ".join(["---"] * len(header)) + " |")
                for row in rows[1:]:
                    # Pad/trim to header length
                    padded = row + [""] * (len(header) - len(row))
                    parts.append("| " + " | ".join(padded[: len(header)]) + " |")
            else:
                parts.append("(empty sheet)")
            parts.append("")

        wb.close()

        text = "\n".join(parts)
        log.info("read_xlsx_done", path=str(path)[-40:], sheets=len(sheets), rows=total_rows)
        return MediaResult(
            success=True,
            text=text,
            metadata={"sheets": len(sheets), "total_rows": total_rows, "format": "xlsx"},
        )

    # ========================================================================
    # Dokument-Analyse (LLM-gestuetzt)
    # ========================================================================

    async def analyze_document(
        self,
        path: str,
        analysis_type: str = "full",
        language: str = "de",
        save_to_vault: bool = False,
    ) -> str:
        """Analysiert ein Dokument strukturiert mit LLM-Unterstuetzung.

        Extrahiert Text, sendet ihn an das LLM mit einem strukturierten
        Analyse-Prompt und liefert eine Analyse mit 6 Abschnitten.

        Args:
            path: Pfad zum Dokument (PDF, DOCX, TXT, HTML, etc.).
            analysis_type: Art der Analyse:
                'full' = vollstaendige Analyse (6 Abschnitte)
                'summary' = nur Zusammenfassung
                'risks' = nur Risiken & Bedenken
                'todos' = nur Handlungsbedarf / To-Dos
            language: Sprache der Analyse ('de' oder 'en').
            save_to_vault: Wenn True, Analyse im Vault speichern.

        Returns:
            Strukturierte Analyse als Markdown-Text.
        """
        if self._llm_fn is None:
            return t("media.analysis_unavailable")

        # 1. Text extrahieren
        extract_result = await self.extract_text(path)
        if not extract_result.success:
            return t("media.analysis_extraction_failed", error=extract_result.error)

        text = extract_result.text
        if not text.strip():
            return t("media.no_text_in_document")

        # Hinweis bei sehr kurzem Text
        ocr_warning = ""
        if len(text.strip()) < 100:
            ocr_warning = t("media.ocr_warning")

        # 2. Analyse-Prompt bauen
        filename = Path(path).name
        prompt = _build_analysis_prompt(text, analysis_type, language, filename)

        # 3. LLM aufrufen
        try:
            analysis = await self._llm_fn(prompt, self._llm_model)
        except Exception as exc:
            log.error("document_analysis_llm_failed", path=path, error=str(exc))
            return t("media.analysis_llm_failed", error=exc)

        result = analysis + ocr_warning

        # 4. Optional im Vault speichern
        if save_to_vault and self._vault is not None:
            try:
                vault_title = f"Analyse: {filename}"
                await self._vault.vault_save(
                    title=vault_title,
                    content=result,
                    tags="analyse, dokument",
                    folder="research",
                    sources=path,
                )
                result += "\n\n" + t("media.vault_saved", title=vault_title)
            except Exception as vault_exc:
                log.warning("vault_save_failed_during_analysis", error=str(vault_exc))
                result += "\n\n" + t("media.vault_save_failed", error=vault_exc)

        log.info("document_analyzed", path=path, type=analysis_type, chars=len(result))
        return result

    # ========================================================================
    # Audio-Konvertierung (ffmpeg)
    # ========================================================================

    # Erlaubte Audio-Samplerates
    ALLOWED_SAMPLE_RATES = frozenset({8000, 11025, 16000, 22050, 32000, 44100, 48000, 96000})

    async def convert_audio(
        self,
        input_path: str,
        output_format: str = "wav",
        *,
        sample_rate: int = 16000,
    ) -> MediaResult:
        """Konvertiert Audio zwischen Formaten via ffmpeg.

        Args:
            input_path: Quell-Audiodatei.
            output_format: Zielformat (wav, mp3, ogg, flac).
            sample_rate: Ziel-Samplerate (8000-96000).
        """
        if sample_rate not in self.ALLOWED_SAMPLE_RATES:
            return MediaResult(
                success=False,
                error=t(
                    "media.invalid_sample_rate",
                    rate=sample_rate,
                    allowed=sorted(self.ALLOWED_SAMPLE_RATES),
                ),
            )

        path = self._validate_input_path(input_path)
        if path is None:
            return MediaResult(
                success=False, error=t("media.file_not_found", path=input_path)
            )

        output_path = self._workspace / f"{path.stem}_converted.{output_format}"

        try:
            proc = await asyncio.create_subprocess_exec(
                "ffmpeg",
                "-i",
                str(path),
                "-ar",
                str(sample_rate),
                "-ac",
                "1",  # Mono
                "-y",  # Überschreiben
                str(output_path),
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
            )
            _, stderr = await proc.communicate()

            if proc.returncode != 0:
                return MediaResult(
                    success=False,
                    error=t("media.ffmpeg_error", error=stderr.decode()[:300]),
                )

            log.info("audio_converted", input=input_path, output=str(output_path))
            return MediaResult(
                success=True,
                text=t("media.audio_converted", path=output_path),
                output_path=str(output_path),
                metadata={"format": output_format, "sample_rate": sample_rate},
            )

        except FileNotFoundError:
            return MediaResult(
                success=False,
                error=t("media.ffmpeg_not_installed"),
            )

    # ========================================================================
    # Bildgroesse aendern (Pillow)
    # ========================================================================

    async def resize_image(
        self,
        image_path: str,
        *,
        max_width: int | None = None,
        max_height: int | None = None,
        output_format: str | None = None,
    ) -> MediaResult:
        """Aendert die Bildgroesse (behaelt Seitenverhaeltnis).

        Args:
            image_path: Quellbild.
            max_width: Maximale Breite (1-max_image_dimension). Default aus Config.
            max_height: Maximale Hoehe (1-max_image_dimension). Default aus Config.
            output_format: Optionales Ausgabeformat (jpg, png, webp).
        """
        if max_width is None:
            max_width = self._default_max_width
        if max_height is None:
            max_height = self._default_max_height
        # Dimensionen validieren
        max_width = max(1, min(max_width, self._max_image_dimension))
        max_height = max(1, min(max_height, self._max_image_dimension))

        path = self._validate_input_path(image_path)
        if path is None:
            return MediaResult(
                success=False, error=t("media.image_not_found", path=image_path)
            )

        try:
            from PIL import Image

            loop = asyncio.get_running_loop()

            def _resize() -> tuple[str, int, int]:
                img = Image.open(path)
                img.thumbnail((max_width, max_height), Image.LANCZOS)

                fmt = output_format or path.suffix.lstrip(".") or "png"
                out = self._workspace / f"{path.stem}_resized.{fmt}"
                img.save(str(out), quality=90)
                return str(out), img.width, img.height

            output, w, h = await loop.run_in_executor(None, _resize)

            log.info("image_resized", input=image_path, output=output, size=f"{w}x{h}")
            return MediaResult(
                success=True,
                text=t("media.image_resized", width=w, height=h, path=output),
                output_path=output,
                metadata={"width": w, "height": h},
            )

        except ImportError:
            return MediaResult(
                success=False,
                error=t("media.pillow_not_installed"),
            )
        except Exception as exc:
            return MediaResult(success=False, error=t("media.image_resize_failed", error=exc))

    # ========================================================================
    # Text → Sprache (TTS)
    # ========================================================================

    # ========================================================================
    # Dokument-Export (PDF / DOCX)
    # ========================================================================

    async def export_document(
        self,
        content: str,
        *,
        fmt: str = "pdf",
        title: str = "",
        author: str = "",
        filename: str = "dokument",
    ) -> MediaResult:
        """Exportiert Text als PDF-, DOCX- oder XLSX-Dokument.

        Args:
            content: Text-Inhalt (Absaetze durch \\n\\n getrennt; fuer xlsx: Zeilen durch \\n,
                Zellen durch | oder Tab getrennt).
            fmt: Ausgabeformat ('pdf', 'docx' oder 'xlsx').
            title: Optionaler Titel/Betreff.
            author: Optionaler Absender/Autor.
            filename: Dateiname ohne Endung.
        """
        fmt = fmt.lower().strip()
        if fmt not in ("pdf", "docx", "xlsx", "pptx"):
            return MediaResult(
                success=False,
                error=t("media.export_unsupported_format", fmt=fmt),
            )

        if not content.strip():
            return MediaResult(success=False, error=t("media.empty_content"))

        # Sicheres Verzeichnis
        doc_dir = Path.home() / ".jarvis" / "workspace" / "documents"
        doc_dir.mkdir(parents=True, exist_ok=True)

        # Dateinamen bereinigen
        safe_name = "".join(c for c in filename if c.isalnum() or c in "-_ ").strip() or "dokument"
        output_path = doc_dir / f"{safe_name}.{fmt}"

        loop = asyncio.get_running_loop()

        try:
            if fmt == "pdf":
                await loop.run_in_executor(
                    None, self._generate_pdf, output_path, content, title, author
                )
            elif fmt == "xlsx":
                await loop.run_in_executor(None, self._generate_xlsx, output_path, content, title)
            elif fmt == "pptx":
                await loop.run_in_executor(
                    None, self._generate_pptx_simple, output_path, content, title
                )
            else:
                await loop.run_in_executor(
                    None, self._generate_docx, output_path, content, title, author
                )

            log.info("document_exported", path=str(output_path), format=fmt)
            return MediaResult(
                success=True,
                text=t("media.document_created", path=output_path),
                output_path=str(output_path),
                metadata={"format": fmt, "title": title, "filename": safe_name},
            )
        except ImportError as exc:
            return MediaResult(success=False, error=str(exc))
        except Exception as exc:
            log.error("document_export_failed", error=str(exc))
            return MediaResult(success=False, error=t("media.document_export_failed", error=exc))

    def _generate_pdf(self, output_path: Path, content: str, title: str, author: str) -> None:
        """Generiert ein PDF-Dokument mit fpdf2."""
        try:
            from fpdf import FPDF
        except ImportError:
            raise ImportError("fpdf2 nicht installiert. pip install fpdf2") from None

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=25)
        pdf.add_page()

        # Unicode-Font einbetten (fuer Umlaute etc.), sonst Helvetica
        font_name = "Helvetica"
        has_bold = True  # Built-in Fonts haben immer Bold
        try:
            # Suche einen Unicode-faehigen TrueType-Font
            font_candidates = [
                # DejaVu Sans (Regular + Bold)
                (
                    "DejaVu",
                    [
                        Path("C:/Windows/Fonts/DejaVuSans.ttf"),
                        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
                        Path.home() / ".fonts" / "DejaVuSans.ttf",
                    ],
                    [
                        Path("C:/Windows/Fonts/DejaVuSans-Bold.ttf"),
                        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"),
                        Path.home() / ".fonts" / "DejaVuSans-Bold.ttf",
                    ],
                ),
                # Arial als Fallback (Windows)
                (
                    "ArialUni",
                    [
                        Path("C:/Windows/Fonts/arial.ttf"),
                    ],
                    [
                        Path("C:/Windows/Fonts/arialbd.ttf"),
                    ],
                ),
            ]
            for fname, regular_paths, bold_paths in font_candidates:
                regular = next((p for p in regular_paths if p.exists()), None)
                if regular:
                    pdf.add_font(fname, "", str(regular), uni=True)
                    font_name = fname
                    bold = next((p for p in bold_paths if p.exists()), None)
                    if bold:
                        pdf.add_font(fname, "B", str(bold), uni=True)
                        has_bold = True
                    else:
                        has_bold = False
                    break
        except Exception:
            pass  # Fallback auf Helvetica

        # Autor/Absender
        if author:
            pdf.set_font(font_name, size=10)
            pdf.multi_cell(0, 6, author)
            pdf.ln(10)

        # Titel
        if title:
            pdf.set_font(font_name, "B" if has_bold else "", size=14)
            pdf.multi_cell(0, 8, title)
            pdf.ln(8)

        # Inhalt
        pdf.set_font(font_name, size=11)
        paragraphs = content.split("\n\n")
        for para in paragraphs:
            para = para.strip()
            if para:
                pdf.multi_cell(0, 6, para)
                pdf.ln(4)

        pdf.output(str(output_path))

    def _generate_docx(self, output_path: Path, content: str, title: str, author: str) -> None:
        """Generiert ein DOCX-Dokument mit python-docx."""
        try:
            from docx import Document
            from docx.shared import Pt
        except ImportError:
            raise ImportError("python-docx nicht installiert. pip install python-docx") from None

        doc = Document()

        # Standardschriftart auf Calibri setzen
        style = doc.styles["Normal"]
        font = style.font
        font.name = "Calibri"
        font.size = Pt(11)

        # Autor/Absender
        if author:
            p = doc.add_paragraph(author)
            p.style.font.size = Pt(10)

        # Titel
        if title:
            doc.add_heading(title, level=1)

        # Inhalt
        paragraphs = content.split("\n\n")
        for para in paragraphs:
            para = para.strip()
            if para:
                doc.add_paragraph(para)

        doc.save(str(output_path))

    def _generate_xlsx(self, output_path: Path, content: str, title: str) -> None:
        """Generate an Excel file from structured text content.

        Content format: rows separated by newlines, cells separated by | or tabs.
        First row is treated as header.
        """
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font
        except ImportError as exc:
            raise ImportError("openpyxl not installed. Run: pip install openpyxl") from exc

        wb = Workbook()
        ws = wb.active
        ws.title = title or "Sheet1"

        for row_idx, line in enumerate(content.strip().split("\n"), 1):
            if not line.strip():
                continue
            # Support | or tab as delimiter
            if "|" in line:
                cells = [c.strip() for c in line.split("|") if c.strip()]
            elif "\t" in line:
                cells = [c.strip() for c in line.split("\t")]
            else:
                cells = [line.strip()]

            for col_idx, value in enumerate(cells, 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=value)
                if row_idx == 1:
                    cell.font = Font(bold=True)

        # Auto-width columns
        for col in ws.columns:
            max_len = 0
            col_letter = col[0].column_letter
            for cell in col:
                if cell.value:
                    max_len = max(max_len, len(str(cell.value)))
            ws.column_dimensions[col_letter].width = min(max_len + 2, 50)

        wb.save(output_path)

    def _generate_pptx_simple(self, output_path: Path, content: str, title: str) -> None:
        """Generate a simple PPTX from plain text (one slide per paragraph)."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError as exc:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx") from exc

        prs = Presentation()
        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title or "Dokument"
        if slide.placeholders[1]:
            slide.placeholders[1].text = ""

        # Content slides (split by double-newline = paragraphs)
        paragraphs = [p.strip() for p in content.split("\n\n") if p.strip()]
        for para in paragraphs:
            slide_layout = prs.slide_layouts[1]  # Title + Content
            slide = prs.slides.add_slide(slide_layout)
            lines = para.split("\n")
            slide.shapes.title.text = lines[0][:80] if lines else ""
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = "\n".join(lines[1:]) if len(lines) > 1 else lines[0]

        prs.save(str(output_path))

    # ========================================================================
    # Strukturiertes Dokument-Erstellen (JSON -> DOCX / PDF / PPTX / XLSX)
    # ========================================================================

    async def create_document(
        self,
        structure: str,
        *,
        fmt: str = "docx",
        filename: str = "dokument",
    ) -> MediaResult:
        """Erstellt ein strukturiertes Dokument aus einer JSON-Eingabe.

        Args:
            structure: JSON-String mit der Dokumentstruktur.
                Format::

                    {
                        "title": "Dokumenttitel",
                        "author": "Autor",
                        "sections": [
                            {"heading": "Abschnitt 1", "content": "Absatztext..."},
                            {"heading": "Abschnitt 2", "content": "Weiterer Text..."},
                            {"table": {"headers": ["Sp1", "Sp2"], "rows": [["a", "b"]]}},
                            {"list": ["Punkt 1", "Punkt 2"]},
                            {"content": "Absatz ohne Ueberschrift"}
                        ]
                    }

            fmt: Ausgabeformat -- 'docx', 'pdf', 'pptx', 'xlsx'.
            filename: Dateiname ohne Endung.
        """
        import json

        fmt = fmt.lower().strip()
        if fmt not in ("docx", "pdf", "pptx", "xlsx"):
            return MediaResult(
                success=False,
                error=t("media.export_unsupported_format", fmt=fmt),
            )

        # JSON parsen
        try:
            doc_struct: dict[str, Any] = json.loads(structure)
        except json.JSONDecodeError as exc:
            return MediaResult(
                success=False,
                error=t("media.invalid_json_variables", error=exc),
            )

        if not isinstance(doc_struct, dict):
            return MediaResult(success=False, error=t("media.invalid_json_variables", error="JSON must be an object"))

        # Ausgabeverzeichnis
        doc_dir = Path.home() / ".jarvis" / "workspace" / "documents"
        doc_dir.mkdir(parents=True, exist_ok=True)

        safe_name = "".join(c for c in filename if c.isalnum() or c in "-_ ").strip() or "dokument"
        output_path = doc_dir / f"{safe_name}.{fmt}"

        loop = asyncio.get_running_loop()

        try:
            if fmt == "docx":
                await loop.run_in_executor(None, self._build_docx, output_path, doc_struct)
            elif fmt == "pdf":
                await loop.run_in_executor(None, self._build_pdf, output_path, doc_struct)
            elif fmt == "pptx":
                await loop.run_in_executor(None, self._build_pptx, output_path, doc_struct)
            else:  # xlsx
                await loop.run_in_executor(None, self._build_xlsx, output_path, doc_struct)

            log.info("document_created", path=str(output_path), format=fmt)
            return MediaResult(
                success=True,
                text=t("media.document_created", path=output_path),
                output_path=str(output_path),
                metadata={
                    "format": fmt,
                    "title": doc_struct.get("title", ""),
                    "filename": safe_name,
                },
            )
        except ImportError as exc:
            return MediaResult(success=False, error=str(exc))
        except Exception as exc:
            log.error("document_create_failed", error=str(exc))
            return MediaResult(success=False, error=t("media.document_create_failed", error=exc))

    def _build_docx(self, path: Path, doc_struct: dict[str, Any]) -> None:
        """Baut ein DOCX-Dokument aus der JSON-Struktur."""
        try:
            from docx import Document
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.shared import Pt
        except ImportError:
            raise ImportError("python-docx nicht installiert. pip install python-docx") from None

        doc = Document()

        title = doc_struct.get("title", "")
        author = doc_struct.get("author", "")

        # Kern-Eigenschaften
        props = doc.core_properties
        if author:
            props.author = author
        if title:
            props.title = title
            doc.add_heading(title, level=0)

        if author:
            p = doc.add_paragraph(author)
            p.alignment = WD_ALIGN_PARAGRAPH.LEFT
            run = p.runs[0] if p.runs else p.add_run(author)
            run.font.size = Pt(10)
            run.font.italic = True

        for section in doc_struct.get("sections", []):
            heading = section.get("heading", "")
            content = section.get("content", "")
            table_data = section.get("table")
            list_items = section.get("list")

            if heading:
                doc.add_heading(heading, level=1)

            if content:
                doc.add_paragraph(content)

            if table_data and isinstance(table_data, dict):
                headers: list[str] = table_data.get("headers", [])
                rows: list[list[str]] = table_data.get("rows", [])
                if headers:
                    tbl = doc.add_table(rows=1 + len(rows), cols=len(headers))
                    tbl.style = "Table Grid"
                    hdr_cells = tbl.rows[0].cells
                    for i, h in enumerate(headers):
                        hdr_cells[i].text = str(h)
                        for run in hdr_cells[i].paragraphs[0].runs:
                            run.font.bold = True
                    for r_idx, row in enumerate(rows, start=1):
                        row_cells = tbl.rows[r_idx].cells
                        for c_idx, cell_val in enumerate(row[: len(headers)]):
                            row_cells[c_idx].text = str(cell_val)

            if list_items and isinstance(list_items, list):
                for item in list_items:
                    doc.add_paragraph(str(item), style="List Bullet")

        doc.save(str(path))

    def _build_pdf(self, path: Path, doc_struct: dict[str, Any]) -> None:
        """Baut ein PDF-Dokument aus der JSON-Struktur."""
        try:
            from fpdf import FPDF
        except ImportError:
            raise ImportError("fpdf2 nicht installiert. pip install fpdf2") from None

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=25)
        pdf.add_page()

        # Font-Ermittlung (Unicode-faehig wenn moeglich)
        font_name = "Helvetica"
        has_bold = True
        try:
            font_candidates = [
                (
                    "DejaVu",
                    [
                        Path("C:/Windows/Fonts/DejaVuSans.ttf"),
                        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
                        Path.home() / ".fonts" / "DejaVuSans.ttf",
                    ],
                    [
                        Path("C:/Windows/Fonts/DejaVuSans-Bold.ttf"),
                        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"),
                        Path.home() / ".fonts" / "DejaVuSans-Bold.ttf",
                    ],
                ),
                (
                    "ArialUni",
                    [Path("C:/Windows/Fonts/arial.ttf")],
                    [Path("C:/Windows/Fonts/arialbd.ttf")],
                ),
            ]
            for fname, regular_paths, bold_paths in font_candidates:
                regular = next((p for p in regular_paths if p.exists()), None)
                if regular:
                    pdf.add_font(fname, "", str(regular), uni=True)
                    font_name = fname
                    bold = next((p for p in bold_paths if p.exists()), None)
                    if bold:
                        pdf.add_font(fname, "B", str(bold), uni=True)
                        has_bold = True
                    else:
                        has_bold = False
                    break
        except Exception:
            pass  # Fallback auf Helvetica

        title = doc_struct.get("title", "")
        author = doc_struct.get("author", "")

        if title:
            pdf.set_font(font_name, "B" if has_bold else "", size=16)
            pdf.multi_cell(0, 10, title)
            pdf.ln(6)

        if author:
            pdf.set_font(font_name, size=10)
            pdf.multi_cell(0, 6, author)
            pdf.ln(8)

        for section in doc_struct.get("sections", []):
            heading = section.get("heading", "")
            content = section.get("content", "")
            table_data = section.get("table")
            list_items = section.get("list")

            if heading:
                pdf.set_font(font_name, "B" if has_bold else "", size=13)
                pdf.multi_cell(0, 8, heading)
                pdf.ln(2)

            if content:
                pdf.set_font(font_name, size=11)
                pdf.multi_cell(0, 6, content)
                pdf.ln(4)

            if table_data and isinstance(table_data, dict):
                headers: list[str] = table_data.get("headers", [])
                rows: list[list[str]] = table_data.get("rows", [])
                if headers:
                    # Effective page width = page width minus left+right margins
                    page_w = pdf.w - pdf.l_margin - pdf.r_margin
                    col_w = min(60.0, page_w / max(len(headers), 1))
                    pdf.set_font(font_name, "B" if has_bold else "", size=9)
                    for h in headers:
                        pdf.cell(col_w, 7, str(h)[:25], border=1)
                    pdf.ln()
                    pdf.set_font(font_name, size=9)
                    for row in rows:
                        for c_idx, cell_val in enumerate(row[: len(headers)]):
                            pdf.cell(col_w, 6, str(cell_val)[:25], border=1)
                        pdf.ln()
                    # Ensure X is reset to left margin after table cells
                    pdf.set_x(pdf.l_margin)
                    pdf.ln(3)

            if list_items and isinstance(list_items, list):
                pdf.set_font(font_name, size=11)
                pdf.set_x(pdf.l_margin)
                for item in list_items:
                    pdf.multi_cell(0, 6, f"- {item}")
                    # multi_cell leaves X at right edge; reset for next item
                    pdf.set_x(pdf.l_margin)
                pdf.ln(3)

        pdf.output(str(path))

    def _build_pptx(self, path: Path, doc_struct: dict[str, Any]) -> None:
        """Baut eine PPTX-Praesentation aus der JSON-Struktur."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            raise ImportError("python-pptx nicht installiert. pip install python-pptx") from None

        prs = Presentation()
        slide_layouts = prs.slide_layouts

        title = doc_struct.get("title", "")
        author = doc_struct.get("author", "")

        # Titelfolie
        title_slide_layout = slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = title
        if slide.placeholders and len(slide.placeholders) > 1:
            slide.placeholders[1].text = author

        for section in doc_struct.get("sections", []):
            heading = section.get("heading", "")
            content = section.get("content", "")
            table_data = section.get("table")
            list_items = section.get("list")

            # Tabellen bekommen eigene Folie
            if table_data and isinstance(table_data, dict):
                headers: list[str] = table_data.get("headers", [])
                rows: list[list[str]] = table_data.get("rows", [])
                if headers:
                    blank_layout = slide_layouts[5]  # Leer
                    tbl_slide = prs.slides.add_slide(blank_layout)
                    if heading:
                        txBox = tbl_slide.shapes.add_textbox(
                            Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
                        )
                        txBox.text_frame.text = heading
                        txBox.text_frame.paragraphs[0].runs[0].font.bold = True
                        txBox.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
                    row_count = 1 + len(rows)
                    col_count = len(headers)
                    tbl = tbl_slide.shapes.add_table(
                        row_count,
                        col_count,
                        Inches(0.5),
                        Inches(1.0),
                        Inches(9.0),
                        Inches(0.4 * row_count),
                    ).table
                    for c_idx, h in enumerate(headers):
                        cell = tbl.cell(0, c_idx)
                        cell.text = str(h)
                        cell.text_frame.paragraphs[0].runs[0].font.bold = True
                    for r_idx, row in enumerate(rows, start=1):
                        for c_idx, cell_val in enumerate(row[:col_count]):
                            tbl.cell(r_idx, c_idx).text = str(cell_val)
                continue  # Naechsten Abschnitt verarbeiten

            # Inhalt-/Listen-Folie
            content_layout = slide_layouts[1]  # Titel + Inhalt
            content_slide = prs.slides.add_slide(content_layout)
            if content_slide.shapes.title:
                content_slide.shapes.title.text = heading or title

            body_placeholder = (
                content_slide.placeholders[1] if len(content_slide.placeholders) > 1 else None
            )
            if body_placeholder:
                tf = body_placeholder.text_frame
                tf.clear()
                if content:
                    tf.text = content
                if list_items and isinstance(list_items, list):
                    for item in list_items:
                        p = tf.add_paragraph()
                        p.text = str(item)
                        p.level = 1

        prs.save(str(path))

    def _build_xlsx(self, path: Path, doc_struct: dict[str, Any]) -> None:
        """Baut eine Excel-Datei aus der JSON-Struktur."""
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Alignment, Font, PatternFill
            from openpyxl.utils import get_column_letter
        except ImportError:
            raise ImportError("openpyxl nicht installiert. pip install openpyxl") from None

        wb = Workbook()
        ws_main = wb.active
        ws_main.title = "Inhalt"

        title = doc_struct.get("title", "")
        author = doc_struct.get("author", "")

        # Titelzeile
        current_row = 1
        if title:
            ws_main.cell(row=current_row, column=1, value=title).font = Font(bold=True, size=14)
            ws_main.merge_cells(
                start_row=current_row, start_column=1, end_row=current_row, end_column=6
            )
            ws_main.cell(row=current_row, column=1).alignment = Alignment(horizontal="center")
            current_row += 1
        if author:
            ws_main.cell(row=current_row, column=1, value=author).font = Font(italic=True, size=10)
            current_row += 1
        if title or author:
            current_row += 1  # Leerzeile

        header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        header_font = Font(bold=True, color="FFFFFF")

        table_sheet_idx = 0
        for section in doc_struct.get("sections", []):
            heading = section.get("heading", "")
            content = section.get("content", "")
            table_data = section.get("table")
            list_items = section.get("list")

            if heading:
                ws_main.cell(row=current_row, column=1, value=heading).font = Font(
                    bold=True, size=12
                )
                current_row += 1

            if content:
                ws_main.cell(row=current_row, column=1, value=content)
                current_row += 1

            if list_items and isinstance(list_items, list):
                for item in list_items:
                    ws_main.cell(row=current_row, column=1, value=f"- {item}")
                    current_row += 1

            if table_data and isinstance(table_data, dict):
                headers: list[str] = table_data.get("headers", [])
                rows: list[list[str]] = table_data.get("rows", [])
                if headers:
                    table_sheet_idx += 1
                    sheet_title = (heading or f"Tabelle {table_sheet_idx}")[:31]
                    ws_tbl = wb.create_sheet(title=sheet_title)
                    # Header
                    for c_idx, h in enumerate(headers, start=1):
                        cell = ws_tbl.cell(row=1, column=c_idx, value=str(h))
                        cell.font = header_font
                        cell.fill = header_fill
                    # Daten
                    for r_idx, row in enumerate(rows, start=2):
                        for c_idx, val in enumerate(row[: len(headers)], start=1):
                            ws_tbl.cell(row=r_idx, column=c_idx, value=str(val))
                    # Spaltenbreite
                    for col_idx, h in enumerate(headers, start=1):
                        col_letter = get_column_letter(col_idx)
                        max_len = len(str(h))
                        for row in rows:
                            if col_idx - 1 < len(row):
                                max_len = max(max_len, len(str(row[col_idx - 1])))
                        ws_tbl.column_dimensions[col_letter].width = min(max_len + 2, 50)
                    # Referenz auf Hauptsheet
                    ws_main.cell(
                        row=current_row,
                        column=1,
                        value=f"[Tabelle: {sheet_title}]",
                    ).font = Font(italic=True, color="4472C4")
                    current_row += 1

            current_row += 1  # Abstand zwischen Abschnitten

        # Spaltenbreite Hauptsheet
        ws_main.column_dimensions["A"].width = 60

        wb.save(str(path))

    # ========================================================================
    # Typst → PDF
    # ========================================================================

    async def typst_render(
        self,
        source: str,
        *,
        filename: str = "dokument",
    ) -> MediaResult:
        """Renders Typst markup to PDF.

        The LLM generates Typst code, this tool compiles it to a high-quality PDF.
        Typst is faster than LaTeX, has cleaner syntax, and produces excellent output.

        Args:
            source: Typst markup code (the content of a .typ file).
            filename: Output filename without extension.

        Returns:
            MediaResult with path to the generated PDF.
        """
        if not source.strip():
            return MediaResult(success=False, error=t("media.empty_typst_source"))

        doc_dir = Path.home() / ".jarvis" / "workspace" / "documents"
        doc_dir.mkdir(parents=True, exist_ok=True)

        safe_name = "".join(c for c in filename if c.isalnum() or c in "-_ ").strip() or "dokument"
        output_path = doc_dir / f"{safe_name}.pdf"

        loop = asyncio.get_running_loop()

        try:
            await loop.run_in_executor(None, self._typst_compile, source, output_path)
            log.info("typst_rendered", path=str(output_path))
            return MediaResult(
                success=True,
                text=t("media.typst_created", path=output_path),
                output_path=str(output_path),
                metadata={"filename": safe_name},
            )
        except ImportError as exc:
            return MediaResult(success=False, error=str(exc))
        except Exception as exc:
            log.error("typst_render_failed", error=str(exc))
            return MediaResult(success=False, error=t("media.typst_failed", error=exc))

    def _typst_compile(self, source: str, output_path: Path) -> None:
        """Sync worker: write .typ temp file, compile, save PDF."""
        import tempfile

        try:
            import typst as typst_lib
        except ImportError as exc:
            raise ImportError("typst not installed. Run: pip install typst") from exc

        tmp = None
        try:
            tmp = tempfile.NamedTemporaryFile(
                suffix=".typ", delete=False, mode="w", encoding="utf-8"
            )
            tmp.write(source)
            tmp.close()  # Must close before typst reads it

            pdf_bytes = typst_lib.compile(tmp.name)
            output_path.write_bytes(pdf_bytes)
        finally:
            if tmp and Path(tmp.name).exists():
                Path(tmp.name).unlink()

    async def text_to_speech(
        self,
        text: str,
        *,
        output_path: str | None = None,
        voice: str = DEFAULT_PIPER_VOICE,
    ) -> MediaResult:
        """Synthetisiert Text zu Sprache (WAV).

        Backend: Piper TTS (lokal, schnell) → eSpeak-NG Fallback

        Args:
            text: Zu sprechender Text.
            output_path: Ausgabedatei. Auto-generiert wenn None.
            voice: Piper-Stimmenmodell.
        """
        if not text.strip():
            return MediaResult(success=False, error=t("media.empty_tts_text"))

        # CWE-22: Validate voice name against path traversal
        from jarvis.security.sanitizer import validate_voice_name

        try:
            validate_voice_name(voice)
        except ValueError as exc:
            return MediaResult(success=False, error=t("media.invalid_voice", error=exc))

        if output_path:
            try:
                out = Path(output_path).expanduser().resolve()
                jarvis_home = Path.home() / ".jarvis"
                try:
                    out.relative_to(self._workspace)
                except ValueError:
                    out.relative_to(jarvis_home)
            except (ValueError, OSError):
                return MediaResult(
                    success=False,
                    error=t("media.output_outside_workspace", path=output_path),
                )
        else:
            out = self._workspace / "tts_output.wav"

        # Versuch 1: Piper
        try:
            proc = await asyncio.create_subprocess_exec(
                "piper",
                "--model",
                voice,
                "--output_file",
                str(out),
                stdin=asyncio.subprocess.PIPE,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
            )
            _, stderr = await proc.communicate(input=text.encode("utf-8"))

            if proc.returncode == 0:
                log.info("tts_piper_success", output=str(out), length=len(text))
                return MediaResult(
                    success=True,
                    text=t("media.tts_audio_created", path=out),
                    output_path=str(out),
                    metadata={"engine": "piper", "voice": voice},
                )
        except FileNotFoundError:
            pass

        # Versuch 2: eSpeak-NG
        try:
            proc = await asyncio.create_subprocess_exec(
                "espeak-ng",
                "-v",
                "de",
                "-w",
                str(out),
                "--",
                text,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
            )
            await proc.communicate()

            if proc.returncode == 0:
                log.info("tts_espeak_success", output=str(out))
                return MediaResult(
                    success=True,
                    text=t("media.tts_audio_created_espeak", path=out),
                    output_path=str(out),
                    metadata={"engine": "espeak"},
                )
        except FileNotFoundError:
            pass

        return MediaResult(
            success=False,
            error=t("media.tts_unavailable"),
        )


# ============================================================================
# Analyse-Prompt-Builder
# ============================================================================


def _build_analysis_prompt(
    text: str,
    analysis_type: str,
    language: str,
    filename: str,
) -> str:
    """Baut den LLM-Prompt fuer die Dokument-Analyse.

    Args:
        text: Extrahierter Dokumenttext.
        analysis_type: full/summary/risks/todos.
        language: Sprache der Analyse.
        filename: Name der Quelldatei.

    Returns:
        Fertiger LLM-Prompt.
    """
    lang_instruction = "Antworte auf Deutsch." if language == "de" else "Answer in English."

    if analysis_type == "summary":
        task = (
            "Erstelle eine prägnante Zusammenfassung des Dokuments in 2-3 Sätzen. "
            "Nenne die wichtigsten Kernaussagen (max. 5 Punkte)."
        )
    elif analysis_type == "risks":
        task = (
            "Identifiziere alle Risiken und Bedenken im Dokument. "
            "Bewerte jedes Risiko als HOCH, MITTEL oder NIEDRIG mit Begründung."
        )
    elif analysis_type == "todos":
        task = (
            "Extrahiere alle Handlungspunkte, To-Dos und offenen Aufgaben aus dem Dokument. "
            "Ordne jedem Punkt eine Priorität zu (Hoch/Mittel/Niedrig)."
        )
    else:  # full
        task = (
            "Erstelle eine vollständige strukturierte Analyse mit genau diesen 6 Abschnitten:\n\n"
            "## Zusammenfassung\n"
            "2-3 Sätze, die den Kern des Dokuments erfassen.\n\n"
            "## Kernaussagen\n"
            "Maximal 7 priorisierte Punkte.\n\n"
            "## Risiken & Bedenken\n"
            "Bewertung: HOCH / MITTEL / NIEDRIG mit Begründung.\n\n"
            "## Handlungsbedarf / To-Dos\n"
            "Konkrete Aktionspunkte mit Priorität.\n\n"
            "## Entscheidungsprotokoll\n"
            "- Bereits getroffene Entscheidungen\n"
            "- Offene Entscheidungen\n\n"
            "## Metadaten\n"
            "- Dokumenttyp (Vertrag/Angebot/Bericht/Protokoll/Rechnung/Sonstiges)\n"
            "- Datum (falls erkennbar)\n"
            "- Beteiligte Parteien\n"
            "- Umfang"
        )

    return (
        f"Du bist ein Experte für Dokumentenanalyse. {lang_instruction}\n\n"
        f"Analysiere das folgende Dokument (Datei: {filename}).\n\n"
        f"AUFGABE:\n{task}\n\n"
        "Falls es sich um ein Rechtsdokument (Vertrag, AGB, etc.) handelt, "
        "füge am Ende diesen Disclaimer an: "
        '"Dies ist eine automatische Analyse und ersetzt keine rechtliche Beratung."\n\n'
        f"DOKUMENT-INHALT:\n---\n{text}\n---"
    )


# ============================================================================
# MCP-Tool-Schemas
# ============================================================================

MEDIA_TOOL_SCHEMAS: dict[str, dict[str, Any]] = {
    "media_transcribe_audio": {
        "description": (
            "Transkribiert eine Audiodatei (WAV, MP3, OGG, etc.) zu Text. "
            "Lokal via Whisper -- keine Cloud-Uploads."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "audio_path": {"type": "string", "description": "Pfad zur Audiodatei"},
                "language": {
                    "type": "string",
                    "description": "Sprache (ISO-Code)",
                    "default": "de",
                },
                "model": {
                    "type": "string",
                    "description": "Whisper-Modell (tiny/base/small/medium/large-v3)",
                    "default": "base",
                },
            },
            "required": ["audio_path"],
        },
    },
    "media_analyze_image": {
        "description": (
            "Analysiert ein Bild mit einem multimodalen Vision-LLM via Ollama. "
            "Beschreibt Inhalt, erkennt Text (OCR), beantwortet Fragen zum Bild. "
            "Standard-Modell: schnell und gut. Mit detail=true: höchste Qualität "
            "(langsamer, für komplexe Bilder, Dokumente, schwer lesbare Texte)."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "image_path": {"type": "string", "description": "Pfad zum Bild"},
                "prompt": {
                    "type": "string",
                    "description": "Analyseanweisung",
                    "default": DEFAULT_IMAGE_PROMPT,
                },
                "detail": {
                    "type": "boolean",
                    "description": "true = Detail-Modell (höchste Qualität, langsamer)",
                    "default": False,
                },
            },
            "required": ["image_path"],
        },
    },
    "media_extract_text": {
        "description": (
            "Extrahiert Text aus Dokumenten: PDF, DOCX, TXT, MD, HTML, CSV, JSON, XML. "
            "Lokal -- keine Cloud-Dienste."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Pfad zum Dokument"},
            },
            "required": ["file_path"],
        },
    },
    "media_convert_audio": {
        "description": "Konvertiert Audio zwischen Formaten (WAV, MP3, OGG, FLAC) via ffmpeg.",
        "inputSchema": {
            "type": "object",
            "properties": {
                "input_path": {"type": "string", "description": "Quell-Audiodatei"},
                "output_format": {
                    "type": "string",
                    "description": "Zielformat",
                    "default": "wav",
                    "enum": ["wav", "mp3", "ogg", "flac"],
                },
                "sample_rate": {
                    "type": "integer",
                    "description": "Samplerate",
                    "default": 16000,
                },
            },
            "required": ["input_path"],
        },
    },
    "media_resize_image": {
        "description": "Ändert Bildgröße (behält Seitenverhältnis). Unterstützt JPG, PNG, WEBP.",
        "inputSchema": {
            "type": "object",
            "properties": {
                "image_path": {"type": "string", "description": "Pfad zum Bild"},
                "max_width": {"type": "integer", "description": "Max. Breite", "default": 1024},
                "max_height": {"type": "integer", "description": "Max. Höhe", "default": 1024},
            },
            "required": ["image_path"],
        },
    },
    "media_tts": {
        "description": ("Konvertiert Text zu Sprache (WAV). Lokal via Piper TTS oder eSpeak-NG."),
        "inputSchema": {
            "type": "object",
            "properties": {
                "text": {"type": "string", "description": "Zu sprechender Text"},
                "voice": {
                    "type": "string",
                    "description": "Piper-Stimmenmodell",
                    "default": DEFAULT_PIPER_VOICE,
                },
            },
            "required": ["text"],
        },
    },
    "analyze_document": {
        "description": (
            "Analysiert ein Dokument (PDF, DOCX, TXT, HTML) strukturiert mit LLM. "
            "Liefert Zusammenfassung, Kernaussagen, Risiken, To-Dos, Entscheidungen und Metadaten. "
            "Optional im Knowledge Vault speichern."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Pfad zum Dokument"},
                "analysis_type": {
                    "type": "string",
                    "enum": ["full", "summary", "risks", "todos"],
                    "description": "Art der Analyse: full (vollständig), summary, risks, todos",
                    "default": "full",
                },
                "language": {
                    "type": "string",
                    "description": "Sprache der Analyse (de/en)",
                    "default": "de",
                },
                "save_to_vault": {
                    "type": "boolean",
                    "description": "Analyse im Knowledge Vault speichern",
                    "default": False,
                },
            },
            "required": ["path"],
        },
    },
    "document_export": {
        "description": (
            "Exportiert Text als PDF-, DOCX-, XLSX- oder PPTX-Dokument. "
            "Fuer PDF/DOCX: Fliesstext, Absaetze durch doppelte Zeilenumbrueche getrennt. "
            "Fuer XLSX: Zeilen durch \\n, Zellen durch | oder Tab getrennt. "
            "Fuer PPTX: Absaetze werden zu einzelnen Slides."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "content": {"type": "string", "description": "Text-Inhalt des Dokuments"},
                "format": {
                    "type": "string",
                    "enum": ["pdf", "docx", "xlsx", "pptx"],
                    "description": "Ausgabeformat",
                    "default": "pdf",
                },
                "title": {"type": "string", "description": "Titel/Betreff", "default": ""},
                "author": {"type": "string", "description": "Absender/Autor", "default": ""},
                "filename": {
                    "type": "string",
                    "description": "Dateiname (ohne Endung)",
                    "default": "dokument",
                },
            },
            "required": ["content"],
        },
    },
    "document_create": {
        "description": (
            "Erstellt ein strukturiertes Dokument (DOCX, PDF, PPTX, XLSX) aus JSON-Struktur. "
            "Unterstuetzt Titel, Abschnitte mit Ueberschriften, Tabellen und Listen. "
            "Ideal fuer Berichte, Praesentationen und strukturierte Dokumente."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "structure": {
                    "type": "string",
                    "description": (
                        'JSON-String: {"title": "...", "author": "...", "sections": ['
                        '{"heading": "...", "content": "..."}, '
                        '{"table": {"headers": [...], "rows": [[...]]}}, '
                        '{"list": ["...", "..."]}'
                        "]}"
                    ),
                },
                "fmt": {
                    "type": "string",
                    "enum": ["docx", "pdf", "pptx", "xlsx"],
                    "description": "Ausgabeformat",
                    "default": "docx",
                },
                "filename": {
                    "type": "string",
                    "description": "Dateiname ohne Endung",
                    "default": "dokument",
                },
            },
            "required": ["structure"],
        },
    },
    "read_pdf": {
        "description": (
            "Liest ein PDF-Dokument strukturiert: Text pro Seite, Metadaten, "
            "optional Bilder und Tabellen extrahieren. Lokal via PyMuPDF."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Pfad zur PDF-Datei"},
                "extract_images": {
                    "type": "boolean",
                    "description": "Eingebettete Bilder extrahieren",
                    "default": False,
                },
                "extract_tables": {
                    "type": "boolean",
                    "description": "Tabellen als Markdown extrahieren",
                    "default": False,
                },
                "pages": {
                    "type": "string",
                    "description": "Seitenbereich (z.B. '1-5', '3', leer=alle)",
                    "default": "",
                },
            },
            "required": ["file_path"],
        },
    },
    "read_ppt": {
        "description": (
            "Liest eine PowerPoint-Praesentation (.pptx): Folieninhalt mit Formatierung, "
            "Sprechernotizen, optional Bilder. Lokal via python-pptx."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Pfad zur PPTX-Datei"},
                "extract_images": {
                    "type": "boolean",
                    "description": "Eingebettete Bilder extrahieren",
                    "default": False,
                },
            },
            "required": ["file_path"],
        },
    },
    "read_docx": {
        "description": (
            "Liest ein Word-Dokument (.docx) strukturiert: Absaetze mit Formatierung, "
            "Tabellen als Markdown, Metadaten, optional Bilder. Lokal via python-docx."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Pfad zur DOCX-Datei"},
                "extract_images": {
                    "type": "boolean",
                    "description": "Eingebettete Bilder extrahieren",
                    "default": False,
                },
                "extract_tables": {
                    "type": "boolean",
                    "description": "Tabellen als Markdown extrahieren",
                    "default": True,
                },
            },
            "required": ["file_path"],
        },
    },
    "read_xlsx": {
        "description": (
            "Liest eine Excel-Datei (.xlsx) strukturiert: Sheets als Markdown-Tabellen, "
            "mit konfigurierbarer Sheet-Auswahl und Zeilenlimit. Lokal via openpyxl."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Pfad zur .xlsx Datei"},
                "sheet_name": {
                    "type": "string",
                    "description": "Spezifisches Sheet (leer = alle Sheets)",
                    "default": "",
                },
                "max_rows": {
                    "type": "integer",
                    "description": "Maximale Zeilen pro Sheet",
                    "default": 500,
                },
            },
            "required": ["file_path"],
        },
    },
    "typst_render": {
        "description": (
            "Kompiliert Typst-Markup zu einem hochwertigen PDF. Typst ist eine moderne "
            "Alternative zu LaTeX mit sauberer Syntax. Der Agent generiert Typst-Code, "
            "dieses Tool kompiliert ihn. Ideal fuer Berichte, Briefe, Rechnungen, "
            "wissenschaftliche Dokumente."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "source": {
                    "type": "string",
                    "description": (
                        "Typst-Markup-Code. Beispiel: "
                        "'= Titel\\n== Abschnitt\\nText...\\n"
                        "#table(columns: 2, [A], [B], [1], [2])'"
                    ),
                },
                "filename": {
                    "type": "string",
                    "description": "Dateiname ohne .pdf Endung",
                    "default": "dokument",
                },
            },
            "required": ["source"],
        },
    },
    "template_list": {
        "description": (
            "Listet alle verfuegbaren Dokumentvorlagen (Templates) auf. "
            "Jede Vorlage hat einen Slug, eine Beschreibung und eine Liste der "
            "auszufuellenden Variablen. Nutze template_render, um eine Vorlage "
            "mit konkreten Werten zu befuellen und als PDF zu exportieren."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {},
            "required": [],
        },
    },
    "template_render": {
        "description": (
            "Befuellt eine Dokumentvorlage mit den angegebenen Variablen und "
            "kompiliert das Ergebnis zu einem PDF. Schritt 1: template_list aufrufen "
            "um verfuegbare Vorlagen zu sehen. Schritt 2: Variablen als JSON-Dict "
            "angeben. Schritt 3: Dieses Tool gibt den Pfad zur fertigen PDF zurueck."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "template_slug": {
                    "type": "string",
                    "description": (
                        "Slug der Vorlage (z.B. 'brief', 'rechnung', 'bericht'). "
                        "Erhaeltlich via template_list."
                    ),
                },
                "variables": {
                    "type": "string",
                    "description": (
                        "JSON-Objekt mit den Platzhalterwerten, z.B. "
                        '\'{"empfaenger_name": "Firma GmbH", "betreff": "Angebot"}\''
                    ),
                    "default": "{}",
                },
                "filename": {
                    "type": "string",
                    "description": "Ausgabe-Dateiname ohne .pdf Endung",
                    "default": "dokument",
                },
            },
            "required": ["template_slug"],
        },
    },
}


def register_media_tools(mcp_client: Any, config: Any = None) -> MediaPipeline:
    """Registriert Media-MCP-Tools beim MCP-Client.

    Args:
        mcp_client: JarvisMCPClient-Instanz.
        config: JarvisConfig-Instanz (optional, fuer Vision-Modell-Auswahl).

    Returns:
        MediaPipeline-Instanz.
    """
    # Vision-Modelle aus Config oder Fallback-Defaults
    vision_model = (
        getattr(config, "vision_model", _DEFAULT_VISION_MODEL) if config else _DEFAULT_VISION_MODEL
    )
    vision_model_detail = (
        getattr(config, "vision_model_detail", _DEFAULT_VISION_MODEL_DETAIL)
        if config
        else _DEFAULT_VISION_MODEL_DETAIL
    )
    ollama_url = (
        getattr(getattr(config, "ollama", None), "base_url", "http://localhost:11434")
        if config
        else "http://localhost:11434"
    )

    pipeline = MediaPipeline(config=config)

    async def _transcribe(
        audio_path: str, language: str = "de", model: str = "base", **_: Any
    ) -> str:
        result = await pipeline.transcribe_audio(
            audio_path,
            language=language,
            model=model,
        )
        return result.text if result.success else t("media.error_detail", error=result.error)

    async def _analyze_image(
        image_path: str, prompt: str = DEFAULT_IMAGE_PROMPT, detail: bool = False, **_: Any
    ) -> str:
        model = vision_model_detail if detail else vision_model
        log.info("vision_model_selected", model=model, detail=detail)
        result = await pipeline.analyze_image(
            image_path,
            prompt=prompt,
            model=model,
            ollama_url=ollama_url,
        )
        return result.text if result.success else t("media.error_detail", error=result.error)

    async def _extract_text(file_path: str, **_: Any) -> str:
        result = await pipeline.extract_text(file_path)
        return result.text if result.success else t("media.error_detail", error=result.error)

    async def _convert_audio(
        input_path: str, output_format: str = "wav", sample_rate: int = 16000, **_: Any
    ) -> str:
        result = await pipeline.convert_audio(
            input_path,
            output_format=output_format,
            sample_rate=sample_rate,
        )
        return result.text if result.success else t("media.error_detail", error=result.error)

    async def _resize_image(
        image_path: str, max_width: int | None = None, max_height: int | None = None, **_: Any
    ) -> str:
        result = await pipeline.resize_image(
            image_path,
            max_width=max_width,
            max_height=max_height,
        )
        return result.text if result.success else t("media.error_detail", error=result.error)

    async def _tts(text: str, voice: str = DEFAULT_PIPER_VOICE, **_: Any) -> str:
        result = await pipeline.text_to_speech(text, voice=voice)
        return result.text if result.success else t("media.error_detail", error=result.error)

    async def _analyze_document(
        path: str,
        analysis_type: str = "full",
        language: str = "de",
        save_to_vault: bool = False,
        **_: Any,
    ) -> str:
        result = await pipeline.analyze_document(
            path,
            analysis_type=analysis_type,
            language=language,
            save_to_vault=save_to_vault,
        )
        return result

    async def _export_document(
        content: str,
        format: str = "pdf",
        title: str = "",
        author: str = "",
        filename: str = "dokument",
        **_: Any,
    ) -> str:
        result = await pipeline.export_document(
            content,
            fmt=format,
            title=title,
            author=author,
            filename=filename,
        )
        if result.success:
            return result.output_path or result.text
        return t("media.error_detail", error=result.error)

    async def _create_document(
        structure: str,
        fmt: str = "docx",
        filename: str = "dokument",
        **_: Any,
    ) -> str:
        result = await pipeline.create_document(
            structure,
            fmt=fmt,
            filename=filename,
        )
        if result.success:
            return result.output_path or result.text
        return t("media.error_detail", error=result.error)

    async def _read_pdf(
        file_path: str,
        extract_images: bool = False,
        extract_tables: bool = False,
        pages: str = "",
        **_: Any,
    ) -> str:
        result = await pipeline.read_pdf(
            file_path,
            extract_images=extract_images,
            extract_tables=extract_tables,
            pages=pages,
        )
        return result.text if result.success else t("media.error_detail", error=result.error)

    async def _read_ppt(
        file_path: str,
        extract_images: bool = False,
        **_: Any,
    ) -> str:
        result = await pipeline.read_ppt(file_path, extract_images=extract_images)
        return result.text if result.success else t("media.error_detail", error=result.error)

    async def _read_docx(
        file_path: str,
        extract_images: bool = False,
        extract_tables: bool = True,
        **_: Any,
    ) -> str:
        result = await pipeline.read_docx(
            file_path,
            extract_images=extract_images,
            extract_tables=extract_tables,
        )
        return result.text if result.success else t("media.error_detail", error=result.error)

    async def _read_xlsx(
        file_path: str,
        sheet_name: str = "",
        max_rows: int = 500,
        **_: Any,
    ) -> str:
        result = await pipeline.read_xlsx(
            file_path,
            sheet_name=sheet_name,
            max_rows=max_rows,
        )
        return result.text if result.success else t("media.error_detail", error=result.error)

    async def _typst_render(
        source: str,
        filename: str = "dokument",
        **_: Any,
    ) -> str:
        result = await pipeline.typst_render(source, filename=filename)
        if result.success:
            return result.output_path or result.text
        return t("media.error_detail", error=result.error)

    def _get_template_manager() -> Any:
        """Lazy-init TemplateManager on first use."""
        if pipeline._template_manager is None:
            from jarvis.documents.templates import TemplateManager

            pipeline._template_manager = TemplateManager()
        return pipeline._template_manager

    async def _template_list(**_: Any) -> str:
        import json as _json

        tm = _get_template_manager()
        templates = tm.list_as_json()
        if not templates:
            return t("media.no_templates")
        return _json.dumps(templates, ensure_ascii=False, indent=2)

    async def _template_render(
        template_slug: str,
        variables: str = "{}",
        filename: str = "dokument",
        **_: Any,
    ) -> str:
        import json as _json

        tm = _get_template_manager()
        try:
            vars_dict: dict[str, str] = _json.loads(variables) if variables.strip() else {}
        except Exception as exc:  # noqa: BLE001
            return t("media.invalid_json_variables", error=exc)

        try:
            rendered_source = tm.render_template(template_slug, vars_dict)
        except KeyError as exc:
            return t("media.template_render_error", error=exc)

        result = await pipeline.typst_render(rendered_source, filename=filename)
        if result.success:
            return result.output_path or result.text
        return t("media.error_detail", error=result.error)

    handlers = {
        "media_transcribe_audio": _transcribe,
        "media_analyze_image": _analyze_image,
        "media_extract_text": _extract_text,
        "media_convert_audio": _convert_audio,
        "media_resize_image": _resize_image,
        "media_tts": _tts,
        "analyze_document": _analyze_document,
        "document_export": _export_document,
        "document_create": _create_document,
        "read_pdf": _read_pdf,
        "read_ppt": _read_ppt,
        "read_docx": _read_docx,
        "read_xlsx": _read_xlsx,
        "typst_render": _typst_render,
        "template_list": _template_list,
        "template_render": _template_render,
    }

    for name, schema in MEDIA_TOOL_SCHEMAS.items():
        mcp_client.register_builtin_handler(
            name,
            handlers[name],
            description=schema["description"],
            input_schema=schema["inputSchema"],
        )

    log.info("media_tools_registered", tools=list(MEDIA_TOOL_SCHEMAS.keys()))
    return pipeline
